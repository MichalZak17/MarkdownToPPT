import re
import tempfile
from pathlib import Path
from dataclasses import dataclass, field
from datetime import datetime
from typing import Optional
from enum import Enum, auto

import typer
from rich.console import Console
from rich.progress import Progress, SpinnerColumn, TextColumn
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image as PILImage

app = typer.Typer(
    name="pptx-gen",
    help="PowerPoint presentation generator from Markdown files",
    add_completion=False,
)
console = Console()


class Theme:
    PRIMARY = RGBColor(0x1E, 0x3A, 0x5F)
    SECONDARY = RGBColor(0x3D, 0x5A, 0x80)
    ACCENT = RGBColor(0xE0, 0x8E, 0x45)

    TEXT_DARK = RGBColor(0x2D, 0x3A, 0x4A)
    TEXT_LIGHT = RGBColor(0xFF, 0xFF, 0xFF)
    TEXT_MUTED = RGBColor(0x6C, 0x75, 0x7D)

    BG_LIGHT = RGBColor(0xF8, 0xF9, 0xFA)
    BG_CODE = RGBColor(0x2D, 0x2D, 0x2D)

    TITLE_SIZE = Pt(36)
    SUBTITLE_SIZE = Pt(20)
    BODY_SIZE = Pt(18)
    CODE_SIZE = Pt(14)
    FOOTER_SIZE = Pt(10)


class SlideType(Enum):
    TITLE = auto()
    SECTION = auto()
    AGENDA = auto()
    CONTENT = auto()
    CODE = auto()
    IMAGE = auto()
    TWO_COLUMN = auto()


@dataclass
class Slide:
    slide_type: SlideType
    title: str = ""
    content: list[str] = field(default_factory=list)
    code: str = ""
    code_language: str = ""
    image_path: str = ""
    notes: str = ""


@dataclass
class Chapter:
    order: int
    filename: str
    title: str
    agenda: str
    slides: list[Slide] = field(default_factory=list)


class MarkdownParser:
    FRONTMATTER_PATTERN = re.compile(r'^---\s*\n(.*?)\n---\s*\n', re.DOTALL)
    SLIDE_SEPARATOR = re.compile(r'\n---\s*\n')
    CODE_BLOCK_PATTERN = re.compile(r'```(\w*)\n(.*?)```', re.DOTALL)
    IMAGE_PATTERN = re.compile(r'!\[([^\]]*)\]\(([^)]+)\)')
    HEADING_PATTERN = re.compile(r'^(#{1,3})\s+(.+)$', re.MULTILINE)
    LIST_ITEM_PATTERN = re.compile(r'^[\s]*[-*+]\s+(.+)$', re.MULTILINE)

    def parse_file(self, filepath: Path) -> Chapter:
        content = filepath.read_text(encoding='utf-8')
        order = self._extract_order(filepath.name)
        title, agenda, content = self._parse_frontmatter(content)

        if not title:
            title = self._filename_to_title(filepath.stem)

        slides = self._parse_slides(content)

        return Chapter(
            order=order,
            filename=filepath.name,
            title=title,
            agenda=agenda,
            slides=slides
        )

    def _extract_order(self, filename: str) -> int:
        match = re.match(r'^(\d+)', filename)
        return int(match.group(1)) if match else 999

    def _filename_to_title(self, stem: str) -> str:
        title = re.sub(r'^\d+[_-]?', '', stem)
        title = title.replace('_', ' ').replace('-', ' ')
        return title.title()

    def _parse_frontmatter(self, content: str) -> tuple[str, str, str]:
        title = ""
        agenda = ""

        match = self.FRONTMATTER_PATTERN.match(content)
        if match:
            frontmatter = match.group(1)
            content = content[match.end():]

            for line in frontmatter.split('\n'):
                if ':' in line:
                    key, value = line.split(':', 1)
                    key = key.strip().lower()
                    value = value.strip().strip('"\'')

                    if key == 'title':
                        title = value
                    elif key == 'agenda':
                        agenda = value

        return title, agenda, content

    def _parse_slides(self, content: str) -> list[Slide]:
        slides = []
        sections = self.SLIDE_SEPARATOR.split(content)

        for section in sections:
            section = section.strip()
            if not section:
                continue

            slide = self._parse_section(section)
            if slide:
                slides.append(slide)

        return slides

    def _parse_section(self, section: str) -> Optional[Slide]:
        code_match = self.CODE_BLOCK_PATTERN.search(section)
        if code_match:
            return self._parse_code_slide(section, code_match)

        image_match = self.IMAGE_PATTERN.search(section)
        if image_match:
            section_without_image = section[:image_match.start()] + section[image_match.end():]
            content_items = self._extract_list_items(section_without_image)
            if content_items:
                title = self._extract_title(section)
                return Slide(
                    slide_type=SlideType.CONTENT,
                    title=title,
                    content=content_items,
                    image_path=image_match.group(2)
                )
            return self._parse_image_slide(section, image_match)

        return self._parse_content_slide(section)

    def _parse_code_slide(self, section: str, code_match: re.Match) -> Slide:
        title = self._extract_title(section)
        language = code_match.group(1) or "text"
        code = code_match.group(2).strip()

        content_before = section[:code_match.start()].strip()
        content_items = self._extract_list_items(content_before)

        return Slide(
            slide_type=SlideType.CODE,
            title=title,
            content=content_items,
            code=code,
            code_language=language
        )

    def _parse_image_slide(self, section: str, image_match: re.Match) -> Slide:
        title = self._extract_title(section)
        image_path = image_match.group(2)

        return Slide(
            slide_type=SlideType.IMAGE,
            title=title,
            image_path=image_path
        )

    def _parse_content_slide(self, section: str) -> Slide:
        title = self._extract_title(section)
        content_items = self._extract_list_items(section)

        if not content_items:
            text = self.HEADING_PATTERN.sub('', section).strip()
            if text:
                content_items = [text]

        return Slide(
            slide_type=SlideType.CONTENT,
            title=title,
            content=content_items
        )

    def _extract_title(self, section: str) -> str:
        match = self.HEADING_PATTERN.search(section)
        return match.group(2).strip() if match else ""

    def _extract_list_items(self, section: str) -> list[str]:
        items = self.LIST_ITEM_PATTERN.findall(section)
        return [item.strip() for item in items]


class PresentationBuilder:
    INLINE_PATTERN = re.compile(
        r'(`[^`]+`)'
        r'|(\*\*[^*]+\*\*)'
        r'|(\*[^*]+\*)'
    )

    def __init__(self, title: str = "Presentation", author: str = "", images_dir: Optional[Path] = None):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.title = title
        self.author = author
        self.images_dir = images_dir or Path("images")
        self.slide_number = 0
        self.total_slides = 0

    def build(self, chapters: list[Chapter], output_path: Path) -> None:
        self.total_slides = 1
        for chapter in chapters:
            self.total_slides += 1
            if chapter.agenda:
                self.total_slides += 1
            self.total_slides += len(chapter.slides)

        self._add_title_slide()

        for chapter in chapters:
            self._add_chapter(chapter)

        self.prs.save(output_path)

    def _add_title_slide(self) -> None:
        self.slide_number += 1
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_slide_background(slide, Theme.PRIMARY)

        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = self.title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = Theme.TEXT_LIGHT
        p.alignment = PP_ALIGN.CENTER

        if self.author:
            author_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.5), Inches(12.333), Inches(0.5)
            )
            tf = author_box.text_frame
            p = tf.paragraphs[0]
            p.text = self.author
            p.font.size = Pt(20)
            p.font.color.rgb = Theme.TEXT_LIGHT
            p.alignment = PP_ALIGN.CENTER

    def _add_chapter(self, chapter: Chapter) -> None:
        self._add_section_slide(chapter)

        if chapter.agenda:
            self._add_agenda_slide(chapter)

        for slide_data in chapter.slides:
            self._add_slide(slide_data, chapter.title)

    def _add_section_slide(self, chapter: Chapter) -> None:
        self.slide_number += 1
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_slide_background(slide, Theme.SECONDARY)

        num_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2), Inches(12.333), Inches(0.8)
        )
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Chapter {chapter.order}"
        p.font.size = Pt(24)
        p.font.color.rgb = Theme.ACCENT
        p.alignment = PP_ALIGN.CENTER

        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.8), Inches(12.333), Inches(1.5)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = chapter.title
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = Theme.TEXT_LIGHT
        p.alignment = PP_ALIGN.CENTER

    def _add_agenda_slide(self, chapter: Chapter) -> None:
        self.slide_number += 1
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_slide_header(slide, "Agenda")

        content_box = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(11.333), Inches(4.5)
        )
        tf = content_box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = chapter.agenda
        p.font.size = Pt(24)
        p.font.color.rgb = Theme.TEXT_DARK
        p.alignment = PP_ALIGN.LEFT

        self._add_slide_footer(slide, chapter.title)

    def _add_slide(self, slide_data: Slide, chapter_title: str) -> None:
        self.slide_number += 1

        if slide_data.slide_type == SlideType.CODE:
            self._add_code_slide(slide_data, chapter_title)
        elif slide_data.slide_type == SlideType.IMAGE:
            self._add_image_slide(slide_data, chapter_title)
        else:
            self._add_content_slide(slide_data, chapter_title)

    def _add_content_slide(self, slide_data: Slide, chapter_title: str) -> None:
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_slide_header(slide, slide_data.title)

        resolved_image = None
        if slide_data.image_path:
            resolved_image = self._resolve_image_path(slide_data.image_path)

        if resolved_image and slide_data.content:
            text_width = Inches(6.5)
        else:
            text_width = Inches(11.333)

        if slide_data.content:
            content_box = slide.shapes.add_textbox(
                Inches(1), Inches(1.8), text_width, Inches(5)
            )
            tf = content_box.text_frame
            tf.word_wrap = True

            for i, item in enumerate(slide_data.content):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()

                self._add_formatted_runs(p, f"\u2022 {item}", Theme.BODY_SIZE, Theme.TEXT_DARK)
                p.space_after = Pt(12)
                p.level = 0

        if resolved_image:
            self._place_image(
                slide, resolved_image,
                Inches(7.8), Inches(2),
                Inches(5), Inches(4.5)
            )

        self._add_slide_footer(slide, chapter_title)

    def _add_formatted_runs(self, paragraph, text: str, size, color) -> None:
        pos = 0
        for match in self.INLINE_PATTERN.finditer(text):
            if match.start() > pos:
                run = paragraph.add_run()
                run.text = text[pos:match.start()]
                run.font.size = size
                run.font.color.rgb = color

            raw = match.group(0)

            if raw.startswith('`'):
                run = paragraph.add_run()
                run.text = raw.strip('`')
                run.font.size = size
                run.font.name = "Consolas"
                run.font.color.rgb = Theme.ACCENT
            elif raw.startswith('**'):
                run = paragraph.add_run()
                run.text = raw.strip('*')
                run.font.size = size
                run.font.color.rgb = color
                run.font.bold = True
            elif raw.startswith('*'):
                run = paragraph.add_run()
                run.text = raw.strip('*')
                run.font.size = size
                run.font.color.rgb = color
                run.font.italic = True

            pos = match.end()

        if pos < len(text):
            run = paragraph.add_run()
            run.text = text[pos:]
            run.font.size = size
            run.font.color.rgb = color

    def _add_code_slide(self, slide_data: Slide, chapter_title: str) -> None:
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_slide_header(slide, slide_data.title)

        code_left = Inches(0.75)
        code_top = Inches(1.8)
        code_width = Inches(11.833)
        code_height = Inches(5)

        code_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            code_left, code_top, code_width, code_height
        )
        code_bg.fill.solid()
        code_bg.fill.fore_color.rgb = Theme.BG_CODE
        code_bg.line.fill.background()

        code_box = slide.shapes.add_textbox(
            code_left + Inches(0.3),
            code_top + Inches(0.3),
            code_width - Inches(0.6),
            code_height - Inches(0.6)
        )
        tf = code_box.text_frame
        tf.word_wrap = False

        lang_label = slide_data.code_language.upper() if slide_data.code_language else "CODE"

        p = tf.paragraphs[0]
        p.text = f"// {lang_label}"
        p.font.size = Pt(12)
        p.font.color.rgb = Theme.TEXT_MUTED
        p.font.name = "Consolas"

        p = tf.add_paragraph()
        p.text = slide_data.code
        p.font.size = Theme.CODE_SIZE
        p.font.color.rgb = Theme.TEXT_LIGHT
        p.font.name = "Consolas"

        self._add_slide_footer(slide, chapter_title)

    def _add_image_slide(self, slide_data: Slide, chapter_title: str) -> None:
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_slide_header(slide, slide_data.title)

        resolved = self._resolve_image_path(slide_data.image_path)
        if resolved:
            self._place_image(
                slide, resolved,
                Inches(2), Inches(2),
                Inches(9.333), Inches(5)
            )
        else:
            placeholder = slide.shapes.add_textbox(
                Inches(2), Inches(3), Inches(9.333), Inches(2)
            )
            tf = placeholder.text_frame
            p = tf.paragraphs[0]
            p.text = f"[Image not found: {slide_data.image_path}]"
            p.font.size = Pt(18)
            p.font.color.rgb = Theme.TEXT_MUTED
            p.alignment = PP_ALIGN.CENTER

        self._add_slide_footer(slide, chapter_title)

    PPTX_SUPPORTED_FORMATS = {".png", ".jpg", ".jpeg", ".gif", ".tiff", ".tif", ".bmp"}

    def _place_image(self, slide, image_path: Path, left, top, max_width, max_height) -> None:
        with PILImage.open(image_path) as img:
            img_w, img_h = img.size

        aspect = img_w / img_h
        width = max_width
        height = int(max_width / aspect)
        if height > max_height:
            height = max_height
            width = int(max_height * aspect)

        if image_path.suffix.lower() not in self.PPTX_SUPPORTED_FORMATS:
            with PILImage.open(image_path) as img:
                img = img.convert("RGBA")
                tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
                img.save(tmp.name, "PNG")
                image_path = Path(tmp.name)

        slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)

    def _resolve_image_path(self, raw_path: str) -> Optional[Path]:
        path = Path(raw_path)
        if path.exists():
            return path
        candidate = self.images_dir / path.name
        if candidate.exists():
            return candidate
        return None

    def _add_slide_header(self, slide, title: str) -> None:
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(13.333), Inches(0.08)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = Theme.ACCENT
        line.line.fill.background()

        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(1)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Theme.TITLE_SIZE
        p.font.bold = True
        p.font.color.rgb = Theme.PRIMARY

    def _add_slide_footer(self, slide, chapter_title: str) -> None:
        num_box = slide.shapes.add_textbox(
            Inches(12.333), Inches(7), Inches(0.75), Inches(0.4)
        )
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{self.slide_number}"
        p.font.size = Theme.FOOTER_SIZE
        p.font.color.rgb = Theme.TEXT_MUTED
        p.alignment = PP_ALIGN.RIGHT

        chapter_box = slide.shapes.add_textbox(
            Inches(0.25), Inches(7), Inches(4), Inches(0.4)
        )
        tf = chapter_box.text_frame
        p = tf.paragraphs[0]
        p.text = chapter_title
        p.font.size = Theme.FOOTER_SIZE
        p.font.color.rgb = Theme.TEXT_MUTED

    def _set_slide_background(self, slide, color: RGBColor) -> None:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color


@app.command()
def build(
    input_dir: Path = typer.Argument(
        Path("markdown"),
        help="Folder containing markdown files",
        exists=True,
        file_okay=False,
        dir_okay=True,
    ),
    output: Optional[Path] = typer.Option(
        None,
        "--output", "-o",
        help="Output file path (defaults to exports/<title>_<timestamp>.pptx)"
    ),
    title: str = typer.Option(
        "Presentation",
        "--title", "-t",
        help="Presentation title"
    ),
    author: str = typer.Option(
        "",
        "--author", "-a",
        help="Presentation author"
    ),
):
    """Compiles Markdown files into a PowerPoint presentation."""
    console.print(f"\n[bold blue]PPTX Presentation Generator[/bold blue]\n")

    md_files = sorted(input_dir.glob("*.md"))

    if not md_files:
        console.print(f"[red]No .md files found in {input_dir}[/red]")
        raise typer.Exit(1)

    if output is None:
        exports_dir = Path("exports")
        exports_dir.mkdir(parents=True, exist_ok=True)
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_title = re.sub(r'[^\w\s-]', '', title).strip().replace(' ', '_').lower()
        output = exports_dir / f"{safe_title}_{timestamp}.pptx"
    else:
        output.parent.mkdir(parents=True, exist_ok=True)

    console.print(f"[dim]Found {len(md_files)} markdown files[/dim]\n")

    parser = MarkdownParser()
    chapters: list[Chapter] = []

    with Progress(
        SpinnerColumn(),
        TextColumn("[progress.description]{task.description}"),
        console=console,
    ) as progress:
        task = progress.add_task("Parsing files...", total=len(md_files))

        for md_file in md_files:
            chapter = parser.parse_file(md_file)
            chapters.append(chapter)
            progress.update(task, advance=1, description=f"Parsing: {md_file.name}")

    chapters.sort(key=lambda c: c.order)

    console.print("\n[bold]Chapters:[/bold]")
    for ch in chapters:
        console.print(f"  {ch.order:02d}. {ch.title} ({len(ch.slides)} slides)")

    console.print(f"\n[dim]Generating presentation...[/dim]")

    images_dir = input_dir / "images"
    builder = PresentationBuilder(title=title, author=author, images_dir=images_dir)
    builder.build(chapters, output)

    console.print(f"\n[green]Saved: {output}[/green]")
    console.print(f"[dim]   Total slides: {builder.slide_number}[/dim]\n")


@app.command()
def preview(
    input_dir: Path = typer.Argument(
        Path("markdown"),
        help="Folder containing markdown files",
        exists=True,
    ),
):
    """Shows a preview of the presentation structure without generating a file."""
    console.print(f"\n[bold blue]Presentation Preview[/bold blue]\n")

    md_files = sorted(input_dir.glob("*.md"))

    if not md_files:
        console.print(f"[red]No .md files found in {input_dir}[/red]")
        raise typer.Exit(1)

    parser = MarkdownParser()
    total_slides = 1

    for md_file in md_files:
        chapter = parser.parse_file(md_file)

        console.print(f"\n[bold cyan]Chapter {chapter.order}: {chapter.title}[/bold cyan]")
        if chapter.agenda:
            console.print(f"  [dim]Agenda: {chapter.agenda[:60]}...[/dim]" if len(chapter.agenda) > 60 else f"  [dim]Agenda: {chapter.agenda}[/dim]")

        total_slides += 1
        if chapter.agenda:
            total_slides += 1

        for slide in chapter.slides:
            type_icon = {
                SlideType.CONTENT: "[cyan]TXT[/cyan]",
                SlideType.CODE: "[green]CODE[/green]",
                SlideType.IMAGE: "[yellow]IMG[/yellow]",
            }.get(slide.slide_type, "[dim]---[/dim]")

            console.print(f"    {type_icon} {slide.title or '(no title)'}")
            total_slides += 1

    console.print(f"\n[bold]Total slides: {total_slides}[/bold]\n")


if __name__ == "__main__":
    app()
